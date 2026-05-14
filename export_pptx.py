"""
v0.4.0 — Export Insight 模組

把 Phase A→D 跑完的成果(query / Q / final_option or final_fig / insight)
打包成一張投影片的 .pptx 報告。

對外只有兩個函式:
    render_chart_to_image(option, Q, chart_engine, fig=None) -> bytes
    build_report_pptx(...) -> bytes

設計原則:
    1. 不依賴 Streamlit;可被 test/CLI 重用。
    2. matplotlib 為主、Plotly 為輔(走 fig.to_image() 若 kaleido 在);全部產 PNG bytes。
    3. ECharts option → matplotlib:萃取 xAxis/yAxis/series,還原成 bar/line/pie/heatmap。
    4. _use_table / _kpi_cards 場景走 python-pptx 原生 table / textbox,不出 chart image。
    5. Insight markdown 解析成段落 + bullet level + bold run,塞進右側 textbox。
"""

from __future__ import annotations

import io
import re
from datetime import datetime, timezone
from typing import Any, Iterable

# matplotlib 必須在 headless backend 下載入(沒 DISPLAY 也能跑)
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as _mpl_fm
import numpy as np
import pandas as pd


# ============================================================
# 🈶 CJK 字體偵測:跨平台優先序(macOS → Linux → Windows → fallback)
# ============================================================
# 重要前提:matplotlib 3.10 對 font.sans-serif 不做 per-glyph fallback —
#   一個 text artist 只會用第一個被找到的字體,因此必須挑「同時涵蓋 Latin + CJK」的字體。
#   下列 _UNIVERSAL_CJK_FONTS 都同時含 Latin glyph(macOS / 主流 CJK 字體都是如此),
#   先試這些;找不到才退回到 Latin-only 字體(此時 CJK 會 tofu,但 ASCII 至少正常)。
_UNIVERSAL_CJK_FONTS = [
    # macOS — 系統字體,Latin + CJK 都涵蓋
    "PingFang TC", "PingFang SC", "PingFang HK",
    "Heiti TC", "Heiti SC", "STHeiti",
    "Hiragino Sans GB", "Hiragino Sans",
    "Apple LiGothic", "Songti TC",
    "Arial Unicode MS",
    # Linux — Noto / Source Han / WenQuanYi 都同時含 Latin
    "Noto Sans CJK TC", "Noto Sans CJK SC", "Noto Sans CJK JP",
    "Noto Sans TC", "Noto Sans SC", "Noto Sans JP",
    "WenQuanYi Zen Hei", "WenQuanYi Micro Hei",
    "Source Han Sans TC", "Source Han Sans SC",
    # Windows
    "Microsoft JhengHei", "Microsoft YaHei",
    "SimHei", "SimSun", "MingLiU",
]
# 純 Latin 字體(只有當系統真的沒裝任何 CJK 字體時才會被選到)
_LATIN_ONLY_FONTS = [
    "Arial", "Helvetica Neue", "Helvetica",
    "DejaVu Sans", "Liberation Sans", "Verdana",
]


def _pick_font_stack() -> list[str]:
    """
    matplotlib font.sans-serif:挑 universal CJK (Latin+CJK) 在前,Latin-only 在後當墊底。
    matplotlib 取第一個存在的字體用到底,所以順序 = 優先序。
    """
    try:
        available = {f.name for f in _mpl_fm.fontManager.ttflist}
    except Exception:
        available = set()
    cjk_picked = [n for n in _UNIVERSAL_CJK_FONTS if n in available]
    latin_picked = [n for n in _LATIN_ONLY_FONTS if n in available]
    stack = cjk_picked + latin_picked
    if not stack:
        stack = ["DejaVu Sans"]
    return stack


# 設一次就好(import 時)
plt.rcParams["font.sans-serif"] = _pick_font_stack()
plt.rcParams["font.family"] = "sans-serif"
plt.rcParams["axes.unicode_minus"] = False  # 避免「-」變方框

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# 🎨 顏色 / 字體 — 與 HR 話圖 Logo 一致
# ============================================================
BRAND_RED = RGBColor(0xD7, 0x19, 0x20)
BRAND_CHARCOAL = RGBColor(0x1F, 0x1F, 0x1F)
BRAND_GRAY_MID = RGBColor(0x77, 0x77, 0x77)
BRAND_GRAY_LIGHT = RGBColor(0xB8, 0xB8, 0xB8)
BRAND_BG_SOFT = RGBColor(0xF7, 0xF7, 0xF7)
BRAND_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# matplotlib 一致的顏色循環(以紅 + 灰階為主)
MPL_COLORS = [
    "#D71920",  # brand red
    "#1F1F1F",  # charcoal
    "#777777",  # gray-mid
    "#4A90D9",  # complementary blue
    "#E08E0B",  # warm amber
    "#52A552",  # green
    "#8E44AD",  # purple
    "#B8B8B8",  # gray-light
]

# 簡報尺寸:16:9 = 13.333" x 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================================================
# 🧰 ECharts option helper:盡量找到 xAxis/yAxis/series 真正內容
# ============================================================
def _option_axis(option: dict, key: str) -> dict:
    """xAxis 或 yAxis 可能是 dict 或 list[dict],統一回第一個 dict。"""
    axis = option.get(key)
    if isinstance(axis, list):
        return axis[0] if axis else {}
    return axis or {}


def _option_axes(option: dict, key: str) -> list:
    """回傳 list[dict];方便處理 dual-axis。"""
    axis = option.get(key)
    if isinstance(axis, list):
        return axis
    if isinstance(axis, dict):
        return [axis]
    return []


def _option_title(option: dict) -> str:
    t = option.get("title")
    if isinstance(t, dict):
        return str(t.get("text") or "")
    if isinstance(t, list) and t:
        return str(t[0].get("text") or "")
    return ""


# ============================================================
# 📈 matplotlib chart renderers
# ============================================================
def _new_fig(width_in: float = 9.0, height_in: float = 5.6):
    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=160)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.tick_params(axis="both", labelsize=10, colors="#1F1F1F")
    ax.grid(axis="y", linestyle="--", alpha=0.4, color="#B8B8B8")
    return fig, ax


def _save_fig_to_bytes(fig) -> bytes:
    buf = io.BytesIO()
    fig.tight_layout()
    fig.savefig(buf, format="png", dpi=160, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _series_iter(option: dict) -> list[dict]:
    series = option.get("series") or []
    return [s for s in series if isinstance(s, dict)]


def _render_bar(option: dict, Q: pd.DataFrame) -> bytes:
    """支援 single / grouped / stacked / 100%-stacked / horizontal。"""
    series = _series_iter(option)
    x_axis = _option_axis(option, "xAxis")
    y_axis = _option_axis(option, "yAxis")

    # 偵測 inverted(橫向 bar)
    x_is_value = (x_axis.get("type") == "value")
    inverted = x_is_value  # xAxis=value & yAxis=category → 橫向

    cat_axis = y_axis if inverted else x_axis
    cats = cat_axis.get("data") or []
    if not cats and not series:
        raise ValueError("Bar render: 找不到分類軸 data 或 series。")

    # 退路:從 series[0].data 推 cats 長度
    if not cats and series:
        data0 = series[0].get("data") or []
        cats = [str(i) for i in range(len(data0))]

    n_cat = len(cats)
    n_series = len(series) or 1

    fig, ax = _new_fig()

    # 是否堆疊?(只要任一 series 有 stack 名稱即視為堆疊)
    stacks = [s.get("stack") for s in series]
    is_stacked = any(stacks) and len({s for s in stacks if s}) <= 2

    bar_width = 0.8 if (is_stacked or n_series == 1) else 0.8 / n_series
    indices = np.arange(n_cat)

    bottoms = np.zeros(n_cat, dtype=float)  # 用於 stacked
    for i, s in enumerate(series):
        raw = s.get("data") or []
        # 補齊長度 / 轉 float
        vals = []
        for v in raw[:n_cat]:
            if isinstance(v, dict):
                vals.append(float(v.get("value", 0) or 0))
            else:
                try:
                    vals.append(float(v))
                except (TypeError, ValueError):
                    vals.append(np.nan)
        while len(vals) < n_cat:
            vals.append(np.nan)
        vals_arr = np.array(vals, dtype=float)

        color = MPL_COLORS[i % len(MPL_COLORS)]
        label = s.get("name") or f"Series {i+1}"

        if is_stacked:
            if inverted:
                ax.barh(indices, vals_arr, left=bottoms,
                        height=0.7, color=color, label=label)
            else:
                ax.bar(indices, vals_arr, bottom=bottoms,
                       width=0.7, color=color, label=label)
            bottoms = bottoms + np.nan_to_num(vals_arr)
        else:
            offset = (i - (n_series - 1) / 2) * bar_width if n_series > 1 else 0
            if inverted:
                ax.barh(indices + offset, vals_arr,
                        height=bar_width, color=color, label=label)
            else:
                ax.bar(indices + offset, vals_arr,
                       width=bar_width, color=color, label=label)

    title = _option_title(option) or "Bar Chart"
    ax.set_title(title, fontsize=14, color="#1F1F1F", fontweight="bold", pad=12)

    cat_label = cat_axis.get("name") or ""
    val_axis = x_axis if inverted else y_axis
    val_label = val_axis.get("name") or ""

    if inverted:
        ax.set_yticks(indices)
        ax.set_yticklabels([str(c) for c in cats], fontsize=10)
        ax.set_xlabel(val_label, fontsize=11)
        ax.set_ylabel(cat_label, fontsize=11)
        # 橫向 bar 的 grid 改 x
        ax.grid(axis="x", linestyle="--", alpha=0.4, color="#B8B8B8")
        ax.grid(axis="y", visible=False)
    else:
        ax.set_xticks(indices)
        ax.set_xticklabels([str(c) for c in cats], fontsize=10,
                           rotation=30 if n_cat > 6 else 0, ha="right" if n_cat > 6 else "center")
        ax.set_xlabel(cat_label, fontsize=11)
        ax.set_ylabel(val_label, fontsize=11)

    if n_series > 1 or any(s.get("name") for s in series):
        ax.legend(loc="best", fontsize=9, frameon=False)

    return _save_fig_to_bytes(fig)


def _render_line(option: dict, Q: pd.DataFrame) -> bytes:
    """支援單軸 / 雙軸 line(yAxisIndex=1 → twinx)。"""
    series = _series_iter(option)
    x_axis = _option_axis(option, "xAxis")
    y_axes = _option_axes(option, "yAxis")

    cats = x_axis.get("data") or []
    if not cats and series:
        data0 = series[0].get("data") or []
        cats = [str(i) for i in range(len(data0))]
    indices = np.arange(len(cats))

    fig, ax = _new_fig()
    ax2 = None
    if len(y_axes) >= 2 and any((s.get("yAxisIndex") or 0) >= 1 for s in series):
        ax2 = ax.twinx()
        ax2.spines["top"].set_visible(False)
        ax2.tick_params(axis="y", labelsize=10, colors="#1F1F1F")

    for i, s in enumerate(series):
        raw = s.get("data") or []
        vals = []
        for v in raw[:len(cats) or len(raw)]:
            if isinstance(v, dict):
                vals.append(float(v.get("value", 0) or 0))
            else:
                try:
                    vals.append(float(v))
                except (TypeError, ValueError):
                    vals.append(np.nan)
        vals_arr = np.array(vals, dtype=float)

        color = MPL_COLORS[i % len(MPL_COLORS)]
        label = s.get("name") or f"Series {i+1}"
        target_ax = ax2 if (ax2 is not None and (s.get("yAxisIndex") or 0) >= 1) else ax

        # 也支援 bar/line 混合(若 type=='bar' 還是有 line 場景常見)
        if s.get("type") == "bar":
            target_ax.bar(indices[:len(vals_arr)], vals_arr,
                          width=0.6, color=color, alpha=0.7, label=label)
        else:
            target_ax.plot(indices[:len(vals_arr)], vals_arr,
                           marker="o", markersize=4, linewidth=2,
                           color=color, label=label)

    title = _option_title(option) or "Line Chart"
    ax.set_title(title, fontsize=14, color="#1F1F1F", fontweight="bold", pad=12)
    ax.set_xticks(indices)
    ax.set_xticklabels([str(c) for c in cats], fontsize=10,
                       rotation=30 if len(cats) > 6 else 0,
                       ha="right" if len(cats) > 6 else "center")
    ax.set_xlabel(x_axis.get("name") or "", fontsize=11)
    ax.set_ylabel(y_axes[0].get("name") if y_axes else "", fontsize=11)
    if ax2 is not None and len(y_axes) >= 2:
        ax2.set_ylabel(y_axes[1].get("name") or "", fontsize=11)

    # 合併雙軸 legend
    handles, labels = ax.get_legend_handles_labels()
    if ax2 is not None:
        h2, l2 = ax2.get_legend_handles_labels()
        handles += h2
        labels += l2
    if handles:
        ax.legend(handles, labels, loc="best", fontsize=9, frameon=False)

    return _save_fig_to_bytes(fig)


def _render_pie(option: dict, Q: pd.DataFrame) -> bytes:
    series = _series_iter(option)
    if not series:
        raise ValueError("Pie render: 沒有 series。")
    data = series[0].get("data") or []
    labels = []
    sizes = []
    for d in data:
        if isinstance(d, dict):
            labels.append(str(d.get("name") or ""))
            try:
                sizes.append(float(d.get("value", 0) or 0))
            except (TypeError, ValueError):
                sizes.append(0.0)
        else:
            labels.append(str(d))
            sizes.append(0.0)
    if not sizes or sum(sizes) == 0:
        raise ValueError("Pie render: 全部 value 為 0,無法畫餅圖。")

    fig, ax = plt.subplots(figsize=(8.5, 5.6), dpi=160)
    fig.patch.set_facecolor("white")
    colors = [MPL_COLORS[i % len(MPL_COLORS)] for i in range(len(sizes))]
    wedges, _texts, autotexts = ax.pie(
        sizes, labels=labels, colors=colors,
        autopct="%1.1f%%", startangle=90,
        wedgeprops={"linewidth": 1.5, "edgecolor": "white"},
        textprops={"fontsize": 10, "color": "#1F1F1F"},
    )
    for t in autotexts:
        t.set_color("white")
        t.set_fontweight("bold")

    title = _option_title(option) or series[0].get("name") or "Pie Chart"
    ax.set_title(title, fontsize=14, color="#1F1F1F", fontweight="bold", pad=12)
    ax.axis("equal")
    return _save_fig_to_bytes(fig)


def _render_heatmap(option: dict, Q: pd.DataFrame) -> bytes:
    series = _series_iter(option)
    if not series:
        raise ValueError("Heatmap render: 沒有 series。")
    x_axis = _option_axis(option, "xAxis")
    y_axis = _option_axis(option, "yAxis")
    xs = x_axis.get("data") or []
    ys = y_axis.get("data") or []
    raw = series[0].get("data") or []
    if not xs or not ys:
        raise ValueError("Heatmap render: xAxis/yAxis.data 不可空。")

    grid = np.full((len(ys), len(xs)), np.nan, dtype=float)
    for item in raw:
        try:
            xi, yi, v = item[0], item[1], item[2]
            grid[int(yi), int(xi)] = float(v)
        except (TypeError, ValueError, IndexError):
            continue

    fig, ax = plt.subplots(figsize=(9.0, 5.6), dpi=160)
    fig.patch.set_facecolor("white")
    im = ax.imshow(grid, cmap="Reds", aspect="auto")
    ax.set_xticks(range(len(xs)))
    ax.set_yticks(range(len(ys)))
    ax.set_xticklabels([str(x) for x in xs], rotation=30, ha="right", fontsize=9)
    ax.set_yticklabels([str(y) for y in ys], fontsize=9)
    cbar = fig.colorbar(im, ax=ax, fraction=0.04, pad=0.04)
    cbar.ax.tick_params(labelsize=8)

    title = _option_title(option) or "Heatmap"
    ax.set_title(title, fontsize=14, color="#1F1F1F", fontweight="bold", pad=12)
    return _save_fig_to_bytes(fig)


# ============================================================
# 🚦 Dispatcher
# ============================================================
def _detect_chart_type(option: dict) -> str:
    """嘗試從 series[0].type 推斷 chart 型別;退路猜 bar。"""
    series = _series_iter(option)
    if not series:
        return "bar"
    t = (series[0].get("type") or "").lower()
    if t in {"bar", "line", "pie", "heatmap", "scatter"}:
        return t
    return "bar"


def render_chart_to_image(
    option: dict | None,
    Q: pd.DataFrame,
    chart_engine: str,
    fig: Any = None,
) -> bytes:
    """主分派:option 轉 png bytes。Plotly 走 fig.to_image。"""
    # Plotly:用 kaleido(若沒裝就 raise,讓 caller 走表格 fallback)
    if chart_engine == "Plotly" and fig is not None:
        return fig.to_image(format="png", width=1600, height=1000, scale=1.5)

    if not isinstance(option, dict):
        raise ValueError("render_chart_to_image: option 不是 dict,無法繪圖。")

    chart_type = _detect_chart_type(option)
    if chart_type == "bar":
        return _render_bar(option, Q)
    if chart_type == "line":
        return _render_line(option, Q)
    if chart_type == "pie":
        return _render_pie(option, Q)
    if chart_type == "heatmap":
        return _render_heatmap(option, Q)
    # scatter / 其他 → 退到 bar(避免炸)
    return _render_bar(option, Q)


# ============================================================
# 📝 Insight markdown → 段落串
# ============================================================
_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
_BULLET_RE = re.compile(r"^(\s*)([-*+•]|\d+\.)\s+(.*)$")
_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")


def _parse_insight_markdown(text: str) -> list[dict]:
    """
    把 Phase D markdown 拆成段落 list:
        [{"kind": "heading", "level": 2, "text": "..."},
         {"kind": "bullet", "level": 0, "runs": [{"text": "...", "bold": True}, ...]},
         {"kind": "para", ...}]
    runs 用來保留 **bold** 內外的差異。
    """
    if not text:
        return []
    lines = [ln.rstrip() for ln in text.splitlines()]
    paragraphs = []
    para_buf = []

    def _flush_para():
        if not para_buf:
            return
        joined = " ".join(s.strip() for s in para_buf if s.strip())
        if joined:
            paragraphs.append({
                "kind": "para",
                "level": 0,
                "runs": _split_bold_runs(joined),
            })
        para_buf.clear()

    for ln in lines:
        if not ln.strip():
            _flush_para()
            continue

        m_head = _HEADING_RE.match(ln.strip())
        if m_head:
            _flush_para()
            level = min(len(m_head.group(1)), 4)
            paragraphs.append({
                "kind": "heading",
                "level": level,
                "text": m_head.group(2).strip(),
                "runs": _split_bold_runs(m_head.group(2).strip()),
            })
            continue

        m_bul = _BULLET_RE.match(ln)
        if m_bul:
            _flush_para()
            indent_len = len(m_bul.group(1))
            level = min(indent_len // 2, 3)  # 每 2 個 space 視為一階
            paragraphs.append({
                "kind": "bullet",
                "level": level,
                "runs": _split_bold_runs(m_bul.group(3).strip()),
            })
            continue

        # 普通段落(累積到下個空行)
        para_buf.append(ln)

    _flush_para()
    return paragraphs


def _split_bold_runs(text: str) -> list[dict]:
    """把 'foo **bar** baz' 拆成 [{'text':'foo ','bold':False},{'text':'bar','bold':True},...]"""
    out = []
    last = 0
    for m in _BOLD_RE.finditer(text):
        if m.start() > last:
            out.append({"text": text[last:m.start()], "bold": False})
        out.append({"text": m.group(1), "bold": True})
        last = m.end()
    if last < len(text):
        out.append({"text": text[last:], "bold": False})
    if not out:
        out.append({"text": text, "bold": False})
    return out


# ============================================================
# 🧱 python-pptx helpers
# ============================================================
def _set_run_style(run, *, size: int = 12, bold: bool = False,
                   color: RGBColor = BRAND_CHARCOAL, font_name: str | None = None):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if font_name:
        run.font.name = font_name


def _add_textbox(slide, left, top, width, height, *, fill: RGBColor | None = None):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(0.1)
    box.text_frame.margin_right = Inches(0.1)
    box.text_frame.margin_top = Inches(0.05)
    box.text_frame.margin_bottom = Inches(0.05)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.fill.background()
    return box


def _set_paragraph_level_indent(paragraph, level: int):
    """python-pptx 對 list level 的 indent 控制是透過 paragraph.level (0~8)。"""
    paragraph.level = max(0, min(level, 4))


def _add_runs_to_paragraph(paragraph, runs: list[dict], *,
                           size: int = 12, color: RGBColor = BRAND_CHARCOAL,
                           font_name: str | None = None):
    """把 run list 寫進 paragraph,保留 bold 區段。"""
    # 第一個 run 直接用 paragraph 自帶
    first = True
    for r in runs:
        if not r.get("text"):
            continue
        if first:
            run = paragraph.add_run()
            first = False
        else:
            run = paragraph.add_run()
        run.text = r["text"]
        _set_run_style(run, size=size, bold=bool(r.get("bold")),
                       color=color, font_name=font_name)


# ============================================================
# 🖼️ 一頁式報告 layout
# ============================================================
def _add_header(slide, query: str, domain: str, source_label: str,
                chart_engine: str, generated_at: datetime):
    # 紅色品牌條
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.18)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_RED
    bar.line.fill.background()

    # 標題 textbox(品牌 + 查詢)
    title_box = _add_textbox(slide, Inches(0.4), Inches(0.28),
                             Inches(12.5), Inches(0.95))
    tf = title_box.text_frame

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    run_hr = p1.add_run()
    run_hr.text = "HR "
    _set_run_style(run_hr, size=24, bold=True, color=BRAND_RED)
    run_ct = p1.add_run()
    run_ct.text = "話圖 "
    _set_run_style(run_ct, size=24, bold=True, color=BRAND_CHARCOAL)
    run_sep = p1.add_run()
    run_sep.text = "·  "
    _set_run_style(run_sep, size=18, color=BRAND_GRAY_MID)
    run_qbody = p1.add_run()
    run_qbody.text = (query or "").strip()[:120]
    _set_run_style(run_qbody, size=16, bold=False, color=BRAND_CHARCOAL)

    # 中繼資料行
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    meta_parts = []
    if domain:
        meta_parts.append(f"領域 · {domain}")
    if source_label:
        meta_parts.append(f"資料 · {source_label}")
    if chart_engine:
        meta_parts.append(f"引擎 · {chart_engine}")
    meta_parts.append(f"生成 · {generated_at.strftime('%Y-%m-%d %H:%M')}")
    run = p2.add_run()
    run.text = "   |   ".join(meta_parts)
    _set_run_style(run, size=10, color=BRAND_GRAY_MID)


def _add_footer(slide, page_label: str = "Generated by HR ChatChart"):
    box = _add_textbox(slide, Inches(0.4), Inches(7.15),
                       Inches(12.5), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = page_label
    _set_run_style(run, size=8, color=BRAND_GRAY_LIGHT)


def _add_chart_picture(slide, image_bytes: bytes, *,
                        left, top, width, height):
    """加 picture,並讓圖維持寬高比置中於 box。"""
    stream = io.BytesIO(image_bytes)
    pic = slide.shapes.add_picture(stream, left, top,
                                    width=width, height=height)
    return pic


def _add_chart_section_title(slide, text: str, left, top, width):
    box = _add_textbox(slide, left, top, width, Inches(0.32))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    _set_run_style(run, size=13, bold=True, color=BRAND_RED)


def _add_table_to_slide(slide, Q: pd.DataFrame, *,
                        left, top, width, height,
                        max_rows: int = 12):
    """把 Q 印成 python-pptx native table。"""
    df = Q.head(max_rows).reset_index(drop=True).copy()
    # 把 NaN 轉空字串避免 'nan' 字樣
    df = df.where(pd.notna(df), "")
    n_rows = len(df) + 1  # +1 表頭
    n_cols = len(df.columns)
    if n_cols == 0:
        # 空表 — 印一段警語
        box = _add_textbox(slide, left, top, width, Inches(0.5))
        p = box.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "(無資料欄位可呈現)"
        _set_run_style(run, size=12, color=BRAND_GRAY_MID)
        return

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table

    # 表頭
    for j, col in enumerate(df.columns):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_CHARCOAL
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(col)
        _set_run_style(run, size=11, bold=True, color=BRAND_WHITE)

    # 資料列
    for i in range(len(df)):
        for j, col in enumerate(df.columns):
            cell = tbl.cell(i + 1, j)
            # 斑馬條
            cell.fill.solid()
            cell.fill.fore_color.rgb = (BRAND_BG_SOFT if i % 2 == 0 else BRAND_WHITE)
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            v = df.iloc[i, j]
            if isinstance(v, float):
                # 簡單 round 處理
                run.text = f"{v:,.2f}" if abs(v) >= 0.01 else f"{v:.4f}"
            else:
                run.text = str(v)
            _set_run_style(run, size=10, color=BRAND_CHARCOAL)


def _add_kpi_cards_to_slide(slide, cards: list[dict], *,
                            left, top, width, height):
    """把 KPI cards 排成 grid(2 欄 / 4 欄取決於數量)。"""
    if not cards:
        return
    n = len(cards)
    if n <= 2:
        cols = n
        rows = 1
    elif n <= 4:
        cols = 2
        rows = 2
    elif n <= 6:
        cols = 3
        rows = 2
    else:
        cols = 4
        rows = (n + 3) // 4

    cell_w = width / cols
    cell_h = height / rows
    gap = Inches(0.12)

    for idx, card in enumerate(cards):
        r = idx // cols
        c = idx % cols
        x = left + c * cell_w + gap
        y = top + r * cell_h + gap
        w = cell_w - 2 * gap
        h = cell_h - 2 * gap

        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BRAND_BG_SOFT
        bg.line.color.rgb = BRAND_GRAY_LIGHT
        bg.line.width = Pt(0.75)
        bg.shadow.inherit = False

        # 顯示 label / value(分兩段)
        tf = bg.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = str(card.get("label") or "")
        _set_run_style(r1, size=11, bold=False, color=BRAND_GRAY_MID)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = str(card.get("value") or "")
        _set_run_style(r2, size=22, bold=True, color=BRAND_RED)

        if card.get("delta"):
            p3 = tf.add_paragraph()
            p3.alignment = PP_ALIGN.LEFT
            r3 = p3.add_run()
            r3.text = str(card["delta"])
            _set_run_style(r3, size=10, bold=False, color=BRAND_GRAY_MID)


def _add_insight_box(slide, insight_text: str, *,
                     left, top, width, height):
    # 標題列
    _add_chart_section_title(slide, "▎  商業洞察",
                              left, top, width)

    body_top = top + Inches(0.4)
    body_h = height - Inches(0.4)

    box = _add_textbox(slide, left, body_top, width, body_h,
                      fill=BRAND_BG_SOFT)
    tf = box.text_frame
    tf.word_wrap = True

    paragraphs = _parse_insight_markdown(insight_text or "")
    if not paragraphs:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "(未啟用 Phase D 商業洞察)"
        _set_run_style(run, size=11, color=BRAND_GRAY_MID)
        return

    # 第一段沿用 paragraphs[0]
    first = True
    for entry in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        kind = entry["kind"]
        if kind == "heading":
            _set_paragraph_level_indent(p, 0)
            _add_runs_to_paragraph(p, entry["runs"],
                                  size=13, color=BRAND_CHARCOAL)
            # 強制把 heading 字體加粗
            for run in p.runs:
                run.font.bold = True
                run.font.color.rgb = BRAND_RED
        elif kind == "bullet":
            _set_paragraph_level_indent(p, entry["level"])
            # 手動加 bullet 點(python-pptx 對 bullet 控制較弱,直接前綴 •)
            prefix_runs = [{"text": "•  ", "bold": False}]
            prefix_runs.extend(entry["runs"])
            _add_runs_to_paragraph(p, prefix_runs,
                                  size=11, color=BRAND_CHARCOAL)
        else:  # para
            _add_runs_to_paragraph(p, entry["runs"],
                                  size=11, color=BRAND_CHARCOAL)
        p.space_after = Pt(3)


# ============================================================
# 📤 對外主函式
# ============================================================
def build_report_pptx(
    *,
    query: str,
    plan_text: str = "",
    Q: pd.DataFrame,
    final_option: dict | None = None,
    final_fig: Any = None,
    insight_text: str | None = None,
    chart_engine: str = "ECharts",
    source_label: str = "",
    domain: str = "",
    use_table_fallback: bool = False,
) -> bytes:
    """
    產出一張投影片的 .pptx(回傳 bytes)。
    版面: 左上 header / 左大塊 chart / 右半 insight / 底部 footer
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # 空白
    slide = prs.slides.add_slide(blank_layout)

    # ----- Header -----
    _add_header(slide, query=query, domain=domain,
                source_label=source_label, chart_engine=chart_engine,
                generated_at=datetime.now(timezone.utc).astimezone())

    # ----- 左半:Chart / Table / KPI -----
    chart_left = Inches(0.4)
    chart_top = Inches(1.6)
    chart_width = Inches(7.6)
    chart_height = Inches(5.3)

    # 場景判斷
    is_kpi = (isinstance(final_option, dict)
              and isinstance(final_option.get("_kpi_cards"), list)
              and final_option["_kpi_cards"])

    if is_kpi:
        _add_chart_section_title(slide, "▎  KPI 卡片",
                                  chart_left, Inches(1.3), chart_width)
        _add_kpi_cards_to_slide(slide, final_option["_kpi_cards"],
                                left=chart_left, top=chart_top,
                                width=chart_width, height=chart_height)
    elif use_table_fallback or (isinstance(final_option, dict)
                                 and final_option.get("_use_table")):
        _add_chart_section_title(slide, "▎  資料表",
                                  chart_left, Inches(1.3), chart_width)
        _add_table_to_slide(slide, Q,
                            left=chart_left, top=chart_top,
                            width=chart_width, height=chart_height)
    else:
        # 嘗試 chart_engine 渲染;失敗就退到表格
        try:
            img_bytes = render_chart_to_image(final_option, Q, chart_engine,
                                              fig=final_fig)
            _add_chart_section_title(slide, "▎  視覺化",
                                      chart_left, Inches(1.3), chart_width)
            # 維持寬高比:用 width 帶入,height 由 picture 自己算
            pic = slide.shapes.add_picture(io.BytesIO(img_bytes),
                                            chart_left, chart_top,
                                            width=chart_width)
            # 太高 → 縮回 chart_height
            if pic.height > chart_height:
                ratio = chart_height / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = chart_height
                # 水平置中
                pic.left = chart_left + (chart_width - pic.width) // 2
        except Exception as exc:
            _add_chart_section_title(
                slide, "▎  資料表(圖形渲染失敗,降級為表格)",
                chart_left, Inches(1.3), chart_width)
            _add_table_to_slide(slide, Q,
                                left=chart_left, top=chart_top,
                                width=chart_width, height=chart_height)
            # 把錯誤摘要塞到 footer 區
            err_box = _add_textbox(slide, chart_left, Inches(6.95),
                                    chart_width, Inches(0.2))
            ep = err_box.text_frame.paragraphs[0]
            er = ep.add_run()
            er.text = f"(圖形渲染失敗:{type(exc).__name__})"
            _set_run_style(er, size=8, color=BRAND_GRAY_MID)

    # ----- 右半:Insight -----
    _add_insight_box(slide, insight_text or "",
                     left=Inches(8.2), top=Inches(1.3),
                     width=Inches(4.75), height=Inches(5.7))

    # ----- Footer -----
    _add_footer(slide)

    # 寫出 bytes
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()
